import io
import random
import re
from collections import Counter

import bleach
import latex2mathml.converter as latex2mathml
import markdown as markdown_lib
import pymupdf as fitz  # `import fitz` is a deprecated alias as of pymupdf 1.28
from pptx import Presentation
from pypdf import PdfReader

from .models import Workspace

# Common English function words + a few chat-filler words, so the dashboard's
# "commonly asked words" view surfaces actual topics instead of "what", "is",
# "please". Not exhaustive — this is deliberately simple keyword frequency
# for the MVP, not real NLP (see CLAUDE.md).
STOPWORDS = frozenset("""
    a about after again against all am an and any are aren't as at be
    because been before being below between both but by can can't cannot
    could couldn't did didn't do does doesn't doing don't down during each
    few for from further had hadn't has hasn't have haven't having he her
    here hers herself him himself his how i if in into is isn't it its
    itself just me more most my myself no nor not now of off on once only
    or other our ours ourselves out over own same she should shouldn't so
    some such than that the their theirs them themselves then there these
    they this those through to too under until up very was wasn't we were
    weren't what when where which while who whom why will with won't would
    wouldn't you your yours yourself yourselves please thanks thank hi
    hello hey ok okay yes ya um uh ve ll re
""".split())

_WORD_RE = re.compile(r"[a-z']+")

# Excludes visually ambiguous characters (0/O, 1/I/L) so codes are easy for
# students to read off a board and type in correctly.
JOIN_CODE_ALPHABET = 'ABCDEFGHJKMNPQRSTUVWXYZ23456789'
JOIN_CODE_LENGTH = 6


def generate_unique_join_code():
    """Return a join code guaranteed not to collide with an existing Workspace."""
    while True:
        code = ''.join(random.choices(JOIN_CODE_ALPHABET, k=JOIN_CODE_LENGTH))
        if not Workspace.objects.filter(join_code=code).exists():
            return code


def extract_pdf_text(content: bytes) -> str:
    """Extract text from PDF bytes, page by page, joined with blank lines.

    Raises pypdf.errors.PyPdfError (or a subclass) if the file can't be
    parsed as a PDF at all — callers should treat that as recoverable (e.g.
    still store the raw file, just without extracted text), since
    scanned/image-only or malformed PDFs are common in the wild.
    """
    reader = PdfReader(io.BytesIO(content))
    return '\n\n'.join(page.extract_text() or '' for page in reader.pages)


def extract_pptx_text(content: bytes) -> str:
    """Extract text from PPTX bytes, slide by slide, joined with blank lines.

    Pulls text from every shape with a text frame (titles, bullets, text
    boxes) on each slide — deliberately excludes speaker notes
    (slide.notes_slide), since those are a teacher's private talking points,
    not something students saw on screen or should have surfaced to them via
    the AI chat's course-material grounding.

    Callers should treat any exception here as recoverable (same as
    extract_pdf_text) — still store the raw file, just without extracted
    text. Unlike pypdf's single PyPdfError, python-pptx doesn't expose one
    reliable exception type for "this isn't a valid .pptx" (a corrupt file
    can surface as a bad zip, a missing part, or an XML parse error
    depending on how it's broken), so callers should catch broadly here.
    """
    presentation = Presentation(io.BytesIO(content))
    slide_texts = []
    for slide in presentation.slides:
        shape_texts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text
        ]
        slide_texts.append('\n'.join(shape_texts))
    return '\n\n'.join(slide_texts)


def rasterize_pdf(content: bytes, zoom: float = 2.0) -> list[bytes]:
    """Render each page of PDF bytes to a PNG, in page order — used by the
    Live Slideshow feature (Lecture Mode only; PDF-only, see Slide model).

    zoom=2.0 is roughly 144 DPI: legible for on-screen presentation without
    bloating storage. Returns one PNG bytes object per page.

    Raises on any failure to open/render — callers should treat this as
    recoverable (same as extract_pdf_text/extract_pptx_text): still store
    the Material, just without Slide rows (no "Present" option for it).
    PyMuPDF doesn't expose one single clean exception type the way pypdf's
    PyPdfError does, so callers should catch broadly here.
    """
    doc = fitz.open(stream=content, filetype='pdf')
    matrix = fitz.Matrix(zoom, zoom)
    return [page.get_pixmap(matrix=matrix).tobytes('png') for page in doc]


def _meaningful_words(text):
    """Lowercase, strip punctuation, and yield the words in `text` that
    aren't stopwords or one/two-letter filler — the shared word-extraction
    rule behind both keyword_frequency (workspace-wide) and
    has_meaningful_content (per-message, used to decide what counts toward
    a student's message count on the teacher dashboard)."""
    for word in _WORD_RE.findall(text.lower()):
        word = word.strip("'")
        if len(word) > 2 and word not in STOPWORDS:
            yield word


def keyword_frequency(text_items, top_n=20):
    """Simple word-frequency count over a batch of message text.

    Per CLAUDE.md: "simple keyword frequency for MVP — skip real
    clustering". Lowercases, strips punctuation, drops stopwords and
    one/two-letter words, and counts what's left.

    Returns a list of (word, count) tuples, most common first.
    """
    counts = Counter()
    for text in text_items:
        counts.update(_meaningful_words(text))
    return counts.most_common(top_n)


def has_meaningful_content(text):
    """True if `text` has at least one word that isn't a stopword/too short
    to be meaningful (same rule as keyword_frequency). Used on the teacher
    dashboard so a message that's just "thanks" or "the" doesn't inflate a
    student's message count."""
    return next(_meaningful_words(text), None) is not None


# --------------------------------------------------------------------------
# AI reply rendering (Markdown + LaTeX)
#
# Gemini's replies (and the Lecture Mode outline generator, same model)
# routinely come back Markdown-formatted — **bold**, bullet lists, headings
# — and, especially for math-heavy Homework/Socratic Mode help, with LaTeX
# math using whichever delimiter convention the model reaches for ($$/$
# or \[\]/\(\)). Left as plain text (the old `|linebreaksbr` treatment),
# all of that shows up as raw literal syntax instead of formatted text —
# this is what render_ai_content fixes. Student-typed messages are
# deliberately NOT run through this (see _message.html) — a student typing
# "it costs $5" shouldn't have that treated as math, and there's no reason
# to markdown-render a student's own plain-text question.
# --------------------------------------------------------------------------

# $$...$$ or \[...\] — display/block math, always unambiguous (no
# legitimate non-math text contains a literal "\[").
_BLOCK_MATH_RE = re.compile(r'\$\$(.+?)\$\$|\\\[(.+?)\\\]', re.DOTALL)

# \(...\) — inline math, also unambiguous.
_BRACKET_INLINE_MATH_RE = re.compile(r'\\\((.+?)\\\)', re.DOTALL)

# $...$ — inline math, but single dollar signs are ambiguous with currency
# ("it costs $5 today and $10 tomorrow"), so this is matched more broadly
# and then filtered by _looks_like_math below before being treated as math.
_DOLLAR_INLINE_MATH_RE = re.compile(r'(?<!\$)\$(?!\s)([^$\n]+?)(?<!\s)\$(?!\$)')

_LOOKS_LIKE_MATH_RE = re.compile(r'[\\^_{}]|[A-Za-z]')
_PURE_NUMBER_RE = re.compile(r'^[\d.,]+$')

_MARKDOWN_EXTENSIONS = ['fenced_code', 'sane_lists', 'nl2br']

# MathML elements latex2mathml can emit, plus the plain-text formatting tags
# Markdown produces from the mode prompts' own encouraged style (headings,
# bullets, bold/italic, code, tables) — anything else (script, img, iframe,
# on* handlers, ...) is stripped. This matters because the AI's reply is
# attacker-adjacent: a student's prior messages can influence what the
# model outputs, so this can't rely on the model "just not" emitting HTML
# (the same reasoning as the keyword-flagging/system-prompt architecture
# elsewhere in this app — see CLAUDE.md).
_ALLOWED_TAGS = [
    'p', 'br', 'strong', 'em', 'ul', 'ol', 'li', 'code', 'pre', 'blockquote', 'hr',
    'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'a', 'div', 'span',
    'table', 'thead', 'tbody', 'tr', 'td', 'th',
    'math', 'mrow', 'msup', 'msub', 'msubsup', 'munder', 'mover', 'munderover',
    'mi', 'mn', 'mo', 'mfrac', 'msqrt', 'mroot', 'mspace', 'mtable', 'mtr', 'mtd',
    'mtext', 'mstyle', 'semantics', 'annotation',
]
_ALLOWED_ATTRS = {
    'a': ['href'],
    'div': ['class'],
    'span': ['class'],
    'code': ['class'],  # fenced_code's "language-xxx" class
    'math': ['xmlns', 'display'],
    'mo': ['stretchy'],
    'mspace': ['width'],
}


def _looks_like_math(content):
    """Heuristic gate for single-$-delimited spans, where $ is genuinely
    ambiguous with currency: skip plain numbers ("$5", "$10.50") and
    require at least one character that wouldn't show up in a dollar
    amount (a backslash, ^, _, {}, or a letter — "x", "\\alpha", "a_1")."""
    stripped = content.strip()
    if not stripped or _PURE_NUMBER_RE.match(stripped):
        return False
    return bool(_LOOKS_LIKE_MATH_RE.search(stripped))


def render_ai_content(text: str) -> str:
    """Render AI-authored text (chat replies, the lecture outline) as safe
    HTML: Markdown formatting plus LaTeX math (converted to MathML, which
    every evergreen browser renders natively — no client-side JS/vendored
    math library needed).

    Math spans are pulled out and replaced with placeholder tokens *before*
    running the Markdown converter, then substituted back in afterward —
    otherwise Markdown's own `_..._` (italic) and `*...*` parsing would
    mangle LaTeX subscripts/multiplication (e.g. "$a_1$", "$\\alpha * \\beta$")
    before latex2mathml ever sees them.

    Returns HTML already passed through bleach.clean() with an explicit
    tag/attribute allowlist — safe to mark as Django `safe` at the call
    site (see templatetags/chat_extras.py), but not a substitute for never
    trusting the *system* prompt (see ai_client.py) — this only protects
    against what the model's *reply* renders as, not what it's told to do.
    """
    placeholders = {}
    counter = 0

    def stash(latex, display):
        nonlocal counter
        counter += 1
        token = f'MATHPLACEHOLDERTOKEN{counter}ENDTOKEN'
        try:
            mathml = latex2mathml.convert(latex.strip(), display=display)
        except Exception:
            # Malformed LaTeX from the model — show the raw source rather
            # than losing the message or crashing the page.
            mathml = bleach.clean(latex.strip())
        tag = 'div' if display == 'block' else 'span'
        css_class = 'ai-math ai-math--block' if display == 'block' else 'ai-math ai-math--inline'
        placeholders[token] = f'<{tag} class="{css_class}">{mathml}</{tag}>'
        return token

    def block_sub(match):
        latex = match.group(1) if match.group(1) is not None else match.group(2)
        # Blank lines around the token so Markdown treats it as its own
        # block rather than folding it into a surrounding paragraph.
        return f'\n\n{stash(latex, "block")}\n\n'

    def bracket_inline_sub(match):
        return stash(match.group(1), 'inline')

    def dollar_inline_sub(match):
        content = match.group(1)
        if not _looks_like_math(content):
            return match.group(0)  # not math-shaped — leave the literal $..$ text alone
        return stash(content, 'inline')

    text = _BLOCK_MATH_RE.sub(block_sub, text)
    text = _BRACKET_INLINE_MATH_RE.sub(bracket_inline_sub, text)
    text = _DOLLAR_INLINE_MATH_RE.sub(dollar_inline_sub, text)

    html = markdown_lib.markdown(text, extensions=_MARKDOWN_EXTENSIONS)
    html = html.replace('<p></p>', '')  # cosmetic: blank paragraphs left around block-math tokens

    for token, replacement in placeholders.items():
        html = html.replace(token, replacement)

    return bleach.clean(html, tags=_ALLOWED_TAGS, attributes=_ALLOWED_ATTRS, strip=True)
